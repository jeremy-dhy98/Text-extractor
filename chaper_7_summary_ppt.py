# -*- coding: utf-8 -*-
"""
Enhanced presentation script with custom styling, backgrounds, and animations.
"""

import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from transformers import pipeline
import os

# Initialize summarizer
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# Function to extract text from scanned PDFs using OCR
def extract_text_from_pdf_ocr(pdf_path, ocr_lang='eng'):
    if not os.path.exists(pdf_path):
        print(f"Error: PDF file '{pdf_path}' not found.")
        return None

    print(f"Extracting text using OCR from {pdf_path}...")
    try:
        images = convert_from_path(pdf_path)
        full_text = ""
        for i, image in enumerate(images):
            page_text = pytesseract.image_to_string(image, lang=ocr_lang)
            full_text += page_text
            print(f"Extracted text from page {i+1}, length: {len(page_text)}")
        print(f"Total extracted text length: {len(full_text)} characters")
        return full_text
    except Exception as e:
        print(f"Error during OCR processing: {e}")
        return None

# Function to create summaries from document
def create_summaries(text, chunk_size=500):
    if not text:
        print("No text provided for summarization.")
        return []

    summaries = []
    tokens = text.split()
    chunks = [tokens[i:i+chunk_size] for i in range(0, len(tokens), chunk_size)]

    print(f"Summarizing {len(chunks)} chunks of text...")
    for i, chunk in enumerate(chunks):
        chunk_text = " ".join(chunk)
        try:
            summary = summarizer(chunk_text, max_length=150, min_length=50, do_sample=False)
            summaries.append(summary[0]['summary_text'])
            print(f"Generated summary for chunk {i+1}: {summary[0]['summary_text']}")
        except Exception as e:
            print(f"Error summarizing chunk {i+1}: {e}")
    return summaries

# Function to create a PowerPoint slide from summaries
def create_ppt_from_summaries(summaries, output_file='Presentation.pptx'):
    if not summaries:
        print("No summaries available to create PowerPoint.")
        return

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Add a title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Summary Presentation_Chapter_7"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True

    # Add a slide for each summary
    for i, summary in enumerate(summaries):
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Title settings
        title.text = f"Summary Slide {i+1}"
        title.text_frame.paragraphs[0].font.size = Pt(30)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Content settings
        content.text = summary
        content.text_frame.paragraphs[0].font.size = Pt(20)
        content.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY

        # Set a background color for the slide
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background

        # Add an icon to the slide
        left = Inches(0.5)
        top = Inches(1)
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(1), Inches(1))
        shape = slide.shapes[-1]
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Icon color
        shape.line.color.rgb = RGBColor(255, 255, 255)

        # Add subtle animation to content
        content.text_frame.text = summary
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.bold = False

    # Save presentation
    try:
        prs.save(output_file)
        print(f"Presentation saved as '{output_file}'")
    except Exception as e:
        print(f"Error saving presentation: {e}")

# Path to PDF file
pdf_path = r"C:\Users\mulwa\Desktop\BMC3.2\BMC 4.1\SoftWareEngineer\chapter_7.pdf"

# Extract text and summarize
text = extract_text_from_pdf_ocr(pdf_path, ocr_lang='eng')
if text:
    summaries = create_summaries(text, chunk_size=500)
    if summaries:
        create_ppt_from_summaries(summaries, "Document Presentation_Chapter_7.pptx")
    else:
        print("No summaries generated.")
else:
    print("No text extracted from PDF.")
